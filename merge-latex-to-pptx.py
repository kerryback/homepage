"""
Merge LaTeX equations from a markdown file into a PowerPoint presentation.

This script parses a markdown file containing LaTeX math equations and inserts
them into specified slides of a PowerPoint deck. Equations are rendered as
high-quality images using matplotlib.

Requirements:
    pip install python-pptx matplotlib pillow

Usage:
    python merge-latex-to-pptx.py input.pptx equations.md output.pptx

Markdown Format:
    The markdown file should specify slide numbers and equations using this format:

    ## Slide 3: Title

    $$x = \frac{-b \pm \sqrt{b^2 - 4ac}}{2a}$$

    Inline math: $E = mc^2$

    $$\int_a^b f(x)\,dx = F(b) - F(a)$$

    The script will extract equations and place them on the specified slide.
"""

import re
import sys
import argparse
from pathlib import Path
from typing import List, Dict, Tuple
import matplotlib.pyplot as plt
import matplotlib
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# Use non-interactive backend
matplotlib.use('Agg')

class LatexEquation:
    """Represents a LaTeX equation with its metadata."""

    def __init__(self, latex: str, is_display: bool, slide_num: int):
        self.latex = latex
        self.is_display = is_display  # True for $$, False for $
        self.slide_num = slide_num

    def __repr__(self):
        eq_type = "display" if self.is_display else "inline"
        return f"Equation(slide={self.slide_num}, type={eq_type}, latex={self.latex[:30]}...)"


def parse_markdown_equations(md_file: Path) -> List[LatexEquation]:
    """
    Parse markdown file and extract LaTeX equations with their slide numbers.

    Args:
        md_file: Path to markdown file

    Returns:
        List of LatexEquation objects
    """
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    equations = []
    current_slide = 1  # Default slide number

    # Split by lines to process sequentially
    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        # Check for slide number specification
        # Format: ## Slide 3 or ## Slide 3: Title
        slide_match = re.match(r'^##\s+Slide\s+(\d+)', line, re.IGNORECASE)
        if slide_match:
            current_slide = int(slide_match.group(1))
            i += 1
            continue

        # Check for display equations ($$...$$)
        if line.strip().startswith('$$'):
            # Multi-line display equation
            equation_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().endswith('$$'):
                equation_lines.append(lines[i])
                i += 1
            if i < len(lines):
                equation_lines.append(lines[i].rstrip('$$').rstrip())

            latex = '\n'.join(equation_lines).strip()
            if latex:
                equations.append(LatexEquation(latex, True, current_slide))
            i += 1
            continue

        # Check for inline equations ($...$)
        inline_matches = re.finditer(r'\$([^\$]+?)\$', line)
        for match in inline_matches:
            latex = match.group(1).strip()
            if latex and not latex.startswith('$'):  # Avoid $$
                equations.append(LatexEquation(latex, False, current_slide))

        i += 1

    return equations


def latex_to_image(latex: str, is_display: bool, dpi: int = 300) -> BytesIO:
    """
    Convert LaTeX equation to PNG image using matplotlib.

    Args:
        latex: LaTeX equation string
        is_display: True for display equations, False for inline
        dpi: Resolution of output image

    Returns:
        BytesIO object containing PNG image data
    """
    # Set font size based on equation type
    fontsize = 24 if is_display else 20

    # Create figure
    fig = plt.figure(figsize=(10, 2))
    fig.patch.set_facecolor('white')
    ax = fig.add_subplot(111)
    ax.axis('off')

    # Wrap in display mode if needed
    if is_display and not latex.startswith(r'\['):
        latex_str = f'${latex}$'
    else:
        latex_str = f'${latex}$'

    # Render equation
    ax.text(0.5, 0.5, latex_str,
            fontsize=fontsize,
            ha='center',
            va='center',
            transform=ax.transAxes)

    # Save to BytesIO
    img_buffer = BytesIO()
    plt.savefig(img_buffer,
                format='png',
                dpi=dpi,
                bbox_inches='tight',
                pad_inches=0.1,
                facecolor='white',
                edgecolor='none')
    plt.close(fig)

    img_buffer.seek(0)
    return img_buffer


def add_equation_to_slide(slide, equation: LatexEquation, position: Tuple[float, float] = None):
    """
    Add a LaTeX equation as an image to a slide.

    Args:
        slide: python-pptx slide object
        equation: LatexEquation object
        position: Optional (left, top) position in inches. If None, centers the equation.
    """
    try:
        # Convert LaTeX to image
        img_buffer = latex_to_image(equation.latex, equation.is_display)

        # Determine position and size
        if position:
            left, top = position
            left = Inches(left)
            top = Inches(top)
        else:
            # Center horizontally, place in lower third of slide
            left = Inches(1.0)
            top = Inches(4.5)

        # Determine width based on equation type
        width = Inches(8.0) if equation.is_display else Inches(4.0)

        # Add picture to slide
        pic = slide.shapes.add_picture(img_buffer, left, top, width=width)

        return pic

    except Exception as e:
        print(f"Error adding equation to slide {equation.slide_num}: {e}")
        print(f"LaTeX: {equation.latex}")
        return None


def merge_equations_to_pptx(pptx_file: Path, md_file: Path, output_file: Path,
                            layout: str = 'auto'):
    """
    Merge LaTeX equations from markdown file into PowerPoint presentation.

    Args:
        pptx_file: Input PowerPoint file
        md_file: Markdown file with equations
        output_file: Output PowerPoint file
        layout: Layout strategy ('auto', 'stacked', 'grid')
    """
    print("="*60)
    print("LaTeX to PowerPoint Merger")
    print("="*60)

    # Parse equations from markdown
    print(f"\nParsing equations from {md_file}...")
    equations = parse_markdown_equations(md_file)
    print(f"Found {len(equations)} equations")

    # Group equations by slide
    equations_by_slide = {}
    for eq in equations:
        if eq.slide_num not in equations_by_slide:
            equations_by_slide[eq.slide_num] = []
        equations_by_slide[eq.slide_num].append(eq)

    print(f"Equations will be added to {len(equations_by_slide)} slides")

    # Load presentation
    print(f"\nLoading presentation from {pptx_file}...")
    prs = Presentation(pptx_file)
    print(f"Presentation has {len(prs.slides)} slides")

    # Add equations to slides
    print("\nAdding equations to slides...")
    for slide_num, slide_equations in sorted(equations_by_slide.items()):
        # Slide numbers are 1-indexed in markdown, but 0-indexed in python-pptx
        slide_idx = slide_num - 1

        if slide_idx < 0 or slide_idx >= len(prs.slides):
            print(f"Warning: Slide {slide_num} not found in presentation (has {len(prs.slides)} slides)")
            continue

        slide = prs.slides[slide_idx]
        print(f"\nSlide {slide_num}: Adding {len(slide_equations)} equation(s)")

        if layout == 'auto':
            # Auto layout: stack equations vertically with spacing
            start_top = 3.5  # Start position
            spacing = 1.5    # Space between equations

            for i, eq in enumerate(slide_equations):
                position = (1.0, start_top + (i * spacing))
                pic = add_equation_to_slide(slide, eq, position)
                if pic:
                    print(f"  ✓ Added: {eq.latex[:50]}...")
                else:
                    print(f"  ✗ Failed: {eq.latex[:50]}...")

        elif layout == 'stacked':
            # Simple stacking with minimal spacing
            start_top = 3.0
            spacing = 1.0

            for i, eq in enumerate(slide_equations):
                position = (1.5, start_top + (i * spacing))
                add_equation_to_slide(slide, eq, position)

        elif layout == 'grid':
            # Grid layout for multiple equations
            cols = 2
            start_left = 1.0
            start_top = 3.0
            col_spacing = 4.5
            row_spacing = 1.5

            for i, eq in enumerate(slide_equations):
                col = i % cols
                row = i // cols
                position = (start_left + (col * col_spacing),
                           start_top + (row * row_spacing))
                add_equation_to_slide(slide, eq, position)

    # Save output
    print(f"\nSaving presentation to {output_file}...")
    prs.save(output_file)

    print("="*60)
    print("✓ Complete!")
    print(f"Output saved to: {output_file}")
    print("="*60)


def create_sample_markdown():
    """Create a sample markdown file demonstrating the format."""
    sample = """## Slide 2: Quadratic Formula

The quadratic formula solves equations of the form $ax^2 + bx + c = 0$:

$$x = \\frac{-b \\pm \\sqrt{b^2 - 4ac}}{2a}$$

## Slide 3: Statistics

The normal distribution with mean $\\mu$ and standard deviation $\\sigma$:

$$f(x) = \\frac{1}{\\sigma\\sqrt{2\\pi}} e^{-\\frac{1}{2}\\left(\\frac{x-\\mu}{\\sigma}\\right)^2}$$

## Slide 5: Calculus

The fundamental theorem of calculus:

$$\\int_a^b f(x)\\,dx = F(b) - F(a)$$

where $F'(x) = f(x)$

## Slide 7: Matrix Operations

$$\\begin{bmatrix}
a_{11} & a_{12} \\\\
a_{21} & a_{22}
\\end{bmatrix}
\\begin{bmatrix}
x_1 \\\\ x_2
\\end{bmatrix} =
\\begin{bmatrix}
b_1 \\\\ b_2
\\end{bmatrix}$$
"""

    with open('sample-equations.md', 'w', encoding='utf-8') as f:
        f.write(sample)

    print("Created sample-equations.md")


def main():
    parser = argparse.ArgumentParser(
        description='Merge LaTeX equations from markdown into PowerPoint',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Merge equations into presentation
  python merge-latex-to-pptx.py input.pptx equations.md output.pptx

  # Use different layout
  python merge-latex-to-pptx.py input.pptx equations.md output.pptx --layout stacked

  # Create sample markdown file
  python merge-latex-to-pptx.py --sample
        """
    )

    parser.add_argument('input_pptx', nargs='?', help='Input PowerPoint file')
    parser.add_argument('equations_md', nargs='?', help='Markdown file with equations')
    parser.add_argument('output_pptx', nargs='?', help='Output PowerPoint file')
    parser.add_argument('--layout', choices=['auto', 'stacked', 'grid'],
                       default='auto', help='Layout strategy for equations')
    parser.add_argument('--sample', action='store_true',
                       help='Create a sample markdown file and exit')

    args = parser.parse_args()

    if args.sample:
        create_sample_markdown()
        return

    if not all([args.input_pptx, args.equations_md, args.output_pptx]):
        parser.print_help()
        sys.exit(1)

    # Convert to Path objects
    pptx_file = Path(args.input_pptx)
    md_file = Path(args.equations_md)
    output_file = Path(args.output_pptx)

    # Validate inputs
    if not pptx_file.exists():
        print(f"Error: Input file not found: {pptx_file}")
        sys.exit(1)

    if not md_file.exists():
        print(f"Error: Markdown file not found: {md_file}")
        sys.exit(1)

    # Merge equations
    merge_equations_to_pptx(pptx_file, md_file, output_file, args.layout)


if __name__ == "__main__":
    main()
